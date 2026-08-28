#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
生成一份「答案已知」的测试 PPT，专门用来验证版式/母版继承那条链路。

为什么必须专门造：
原 PPT 工具的缺陷是「整层看不见」—— 版式和母版上的元素一个都不输出，
而输出本身看起来毫无异常（文字都在，只是页面里少了 logo 和色条）。
拿真实 PPT 测只能看出「好像少了点东西」，说不清少了什么。
造一份自己知道每个数值的文件，才能逐项断言。

文件里埋的东西（测试脚本按这些断言）：
  母版上   : 一个 #1F6F8B 的顶部色条（矩形，名 MasterBar）
  版式上   : 一个 #C6402E 的角标（矩形，名 LayoutBadge）
  母版背景 : #F6F8FA 纯色
  第 1 页  : 标题占位符，只写文字不写位置 ⇒ 位置必须从版式继承
             另有一个写死位置的普通矩形（x=1cm y=5cm 3×2cm，填充 #1E7A46）
  第 2 页  : 一个 hidden=1 的形状（必须被跳过并计数）
             一个 2×3 表格
用法： python make_pptx.py <输出路径>
"""
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn


def solid(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(rgb)
    shape.line.fill.background()


def rect_sp(sid, name, x, y, cx, cy, rgb):
    """直接造一个纯色矩形的 <p:sp>。

    python-pptx 的 MasterShapes / LayoutShapes 没有 add_shape，
    只能往 spTree 里塞 XML —— 而我们要测的恰恰是母版和版式上的形状。
    """
    xml = (
        '<p:sp %s>'
        '<p:nvSpPr><p:cNvPr id="%d" name="%s"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        '<p:spPr>'
        '<a:xfrm><a:off x="%d" y="%d"/><a:ext cx="%d" cy="%d"/></a:xfrm>'
        '<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
        '<a:solidFill><a:srgbClr val="%s"/></a:solidFill>'
        '<a:ln><a:noFill/></a:ln>'
        '</p:spPr>'
        '<p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody>'
        '</p:sp>'
    ) % (nsdecls("p", "a"), sid, name, int(x), int(y), int(cx), int(cy), rgb)
    return parse_xml(xml)


def main(dest: Path) -> int:
    prs = Presentation()
    prs.slide_width = Cm(33.867)      # 16:9
    prs.slide_height = Cm(19.05)

    master = prs.slide_master

    # ---- 母版背景（python-pptx 没有高层 API，直接改 XML）----
    cSld = master.element.find(qn("p:cSld"))
    bg = cSld.makeelement(qn("p:bg"), {})
    bgPr = bg.makeelement(qn("p:bgPr"), {})
    solidFill = bgPr.makeelement(qn("a:solidFill"), {})
    srgb = solidFill.makeelement(qn("a:srgbClr"), {"val": "F6F8FA"})
    solidFill.append(srgb)
    bgPr.append(solidFill)
    bgPr.append(bgPr.makeelement(qn("a:effectLst"), {}))
    bg.append(bgPr)
    cSld.insert(0, bg)

    # ---- 母版上的顶部色条：原工具完全看不见的那一层 ----
    master.shapes._spTree.append(
        rect_sp(900, "MasterBar", 0, 0, prs.slide_width, Cm(1.2), "1F6F8B"))

    layout = prs.slide_layouts[1]     # 标题和内容
    layout.shapes._spTree.append(
        rect_sp(901, "LayoutBadge", Cm(30.0), Cm(0.3), Cm(3.0), Cm(0.6), "C6402E"))

    # ---- 第 1 页：标题只写文字，不碰位置 ⇒ 必须从版式继承 ----
    s1 = prs.slides.add_slide(layout)
    s1.shapes.title.text = "继承测试标题"
    body = s1.placeholders[1]
    body.text_frame.text = "第一行正文"
    body.text_frame.add_paragraph().text = "第二行正文"

    box = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(1), Cm(5), Cm(3), Cm(2))
    box.name = "ExplicitBox"
    solid(box, "1E7A46")

    # ---- 第 2 页：隐藏形状 + 表格 ----
    s2 = prs.slides.add_slide(prs.slide_layouts[5])   # 仅标题
    s2.shapes.title.text = "第二页"

    ghost = s2.shapes.add_shape(MSO_SHAPE.OVAL, Cm(2), Cm(2), Cm(2), Cm(2))
    ghost.name = "GhostShape"
    solid(ghost, "FF0000")
    ghost.element.nvSpPr.cNvPr.set("hidden", "1")

    tbl = s2.shapes.add_table(2, 3, Cm(2), Cm(8), Cm(18), Cm(3)).table
    cells = [["Brand", "Units", "Share"], ["Samsung", "4570k", "31%"]]
    for r, row in enumerate(cells):
        for c, t in enumerate(row):
            tbl.cell(r, c).text = t

    dest.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(dest))
    print(f"已生成 {dest}  ({dest.stat().st_size/1024:.0f} KB)")
    return 0


if __name__ == "__main__":
    out = Path(sys.argv[1] if len(sys.argv) > 1 else "test-inherit.pptx")
    raise SystemExit(main(out))
